from flask import Flask, render_template, request, redirect, url_for
from werkzeug.utils import secure_filename
from pptx import Presentation
from gtts import gTTS
import os
import ffmpeg
import random
import math

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
STATIC_FOLDER = 'static'
VIDEO_OUTPUT_FOLDER = os.path.join(STATIC_FOLDER, 'videos')
BACKGROUND_FOLDER = os.path.join(STATIC_FOLDER, 'backgrounds')
TEMP_FOLDER = os.path.join(STATIC_FOLDER, 'temp')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(VIDEO_OUTPUT_FOLDER, exist_ok=True)
os.makedirs(BACKGROUND_FOLDER, exist_ok=True)
os.makedirs(TEMP_FOLDER, exist_ok=True)

# Function to extract text from PPT
def extract_text_from_ppt(ppt_file_path):
    try:
        prs = Presentation(ppt_file_path)
        slides_text = []
        for slide in prs.slides:
            text = ''
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + ' '
            if text.strip():
                slides_text.append(text.strip())
        return slides_text
    except Exception as e:
        print(f"Error processing PPT: {e}")
        return []

# Function to convert text to audio
def text_to_audio(text, output_filename, lang='en'):
    try:
        tts = gTTS(text=text, lang=lang)
        tts.save(output_filename)
        return output_filename
    except Exception as e:
        print(f"Error generating audio: {e}")
        return None

# Function to generate video with fading backgrounds and selected character
def generate_static_video(audio_path, output_path, character):
    image_path = os.path.join(STATIC_FOLDER, 'img', f'{character}.png')
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Character image not found at {image_path}")
    print(f"Using character image: {image_path}")

    # List background images
    background_files = [f for f in os.listdir(BACKGROUND_FOLDER) if f.endswith(('.jpg', '.png'))]
    if not background_files:
        raise FileNotFoundError(f"No background images found in {BACKGROUND_FOLDER}")
    print(f"Found {len(background_files)} background images: {background_files}")

    # Audio duration
    audio_duration = float(ffmpeg.probe(audio_path)['format']['duration'])
    bg_change_duration = 5  # Change background every 5 seconds
    num_bg_segments = math.ceil(audio_duration / bg_change_duration)
    print(f"Audio duration: {audio_duration}s, Segments: {num_bg_segments}")

    # Cycle through backgrounds
    background_cycle = (background_files * (num_bg_segments // len(background_files) + 1))[:num_bg_segments]
    random.shuffle(background_cycle)

    # Generate segments with fading backgrounds
    segment_files = []
    for i in range(num_bg_segments):
        start_time = i * bg_change_duration
        duration = min(bg_change_duration, audio_duration - start_time)
        if duration <= 0:
            break

        bg_file = background_cycle[i]
        bg_path = os.path.join(BACKGROUND_FOLDER, bg_file)
        segment_video = os.path.join(TEMP_FOLDER, f'segment_{i}.mp4')
        segment_audio = os.path.join(TEMP_FOLDER, f'segment_audio_{i}.mp3')

        # Extract audio segment
        try:
            audio = ffmpeg.input(audio_path, ss=start_time, t=duration)
            ffmpeg.output(audio, segment_audio, acodec='mp3', loglevel='quiet').run(overwrite_output=True)
            print(f"Generated audio segment: {segment_audio}")
        except ffmpeg.Error as e:
            print(f"Error generating audio segment {i}: {e.stderr.decode() if e.stderr else 'No stderr output'}")
            raise

        # Background and character setup
        try:
            bg_stream = ffmpeg.input(bg_path, loop=1, framerate=24, t=duration)

            # Apply fade effect
            if i > 0 and i < num_bg_segments - 1:
                bg_stream = bg_stream.filter('fade', type='in', duration=1, start_time=0)
                bg_stream = bg_stream.filter('fade', type='out', duration=1, start_time=duration-1)
            elif i > 0:
                bg_stream = bg_stream.filter('fade', type='in', duration=1, start_time=0)
            elif i < num_bg_segments - 1:
                bg_stream = bg_stream.filter('fade', type='out', duration=1, start_time=duration-1)

            char_stream = ffmpeg.input(image_path, loop=1, framerate=24, t=duration)
            char_stream = char_stream.filter('scale', 'iw*0.75', 'ih*0.75')

            # Overlay character at center
            video = ffmpeg.overlay(bg_stream, char_stream, 
                                   x='(W-w)/2', y='(H-h)/2', shortest=1)

            # Output segment
            audio_segment = ffmpeg.input(segment_audio)
            output = ffmpeg.output(
                video, audio_segment, segment_video,
                vcodec='libx264', acodec='aac', pix_fmt='yuv420p',
                shortest=None
            )
            ffmpeg.run(output, overwrite_output=True)
            print(f"Generated video segment: {segment_video}")
        except ffmpeg.Error as e:
            print(f"Error generating video segment {i}: {e.stderr.decode() if e.stderr else 'No stderr output'}")
            raise

        segment_files.append(os.path.abspath(segment_video))

    # Create concat list file
    concat_list_path = os.path.join(TEMP_FOLDER, 'concat_list.txt')
    with open(concat_list_path, 'w') as f:
        for segment in segment_files:
            f.write(f"file '{segment}'\n")
    print(f"Created concat list: {concat_list_path}")

    # Concatenate segments
    try:
        concat_stream = ffmpeg.input(concat_list_path, format='concat', safe=0)
        audio_full = ffmpeg.input(audio_path)
        final_output = ffmpeg.output(
            concat_stream, audio_full, output_path,
            c='copy', acodec='aac', shortest=None
        )
        ffmpeg.run(final_output, overwrite_output=True)
        print(f"Final video generated: {output_path}")
    except ffmpeg.Error as e:
        print(f"Error during concatenation: {e.stderr.decode() if e.stderr else 'No stderr output'}")
        raise

    # Clean up
    for f in segment_files + [os.path.join(TEMP_FOLDER, f'segment_audio_{i}.mp3') for i in range(len(segment_files))] + [concat_list_path]:
        if os.path.exists(f):
            os.remove(f)
            print(f"Cleaned up: {f}")

# Routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/login')
def login():
    return render_template('login.html', dashboard_url=url_for('dashboard'))

@app.route('/dashboard')
def dashboard():
    return render_template('dashboard.html')

@app.route('/upload_selection')
def upload_selection():
    return render_template('upload_selection.html')

@app.route('/logout')
def logout():
    return redirect(url_for('index'))

@app.route('/video')
def video():
    return render_template('video.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return "No file part in request", 400
        file = request.files['file']
        if file.filename == '':
            return "No file selected", 400
        if not file.filename.endswith('.pptx'):
            return "Invalid file format. Please upload a .pptx file", 400

        character = request.form.get('character')
        if not character:
            return "No character selected", 400

        filename = secure_filename(file.filename)
        ppt_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(ppt_path)

        slides_text = extract_text_from_ppt(ppt_path)
        if not slides_text:
            return "No text extracted from PPT", 400
        full_text = ' '.join(slides_text)

        audio_path = os.path.join(UPLOAD_FOLDER, 'audio.mp3')
        if not text_to_audio(full_text, audio_path):
            return "Failed to generate audio", 500

        video_path = os.path.join(VIDEO_OUTPUT_FOLDER, 'generated_video.mp4')
        generate_static_video(audio_path, video_path, character)

        return redirect(url_for('video'))
    except Exception as e:
        print(f"Upload error: {str(e)}")
        return f"Error processing upload: {str(e)}", 500

if __name__ == '__main__':
    app.run(debug=True)